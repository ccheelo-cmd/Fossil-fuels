"""Build the v2 'Case For Fossil Fuels' .pptx — reliability spine.

This is a NEW deck; it does NOT overwrite the original build_deck.py output.
Output: deliverable/case-for-fossil-fuels-v2.pptx

Narrative spine (Storytelling-with-Data, SCR):
  Situation  — reliable energy is survival; fossil fuels are ~80% of it.
  Complication — clean energy alone can't guarantee that reliability yet:
       Zambia's drought (South) + Alberta's cold snap (North) + the justice
       asymmetry + the energy->life curve.
  Resolution — a managed transition: gas bridge -> renewables stacked ->
       retire firm power on an engineering schedule.

Two-accent palette used thematically:
  BLUE (#1F4E79) = path forward / what's built;  AMBER (#D97706) = cost / stakes.
9 slides. All citations consolidated on the final Sources slide.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Palette (thematic two-accent) ----
BLUE        = RGBColor(0x1F, 0x4E, 0x79)
BLUE_LIGHT  = RGBColor(0xDB, 0xE5, 0xF1)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_LIGHT = RGBColor(0xFD, 0xE8, 0xCC)
TEXT        = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED  = RGBColor(0x4B, 0x55, 0x63)
GREY_CTX    = RGBColor(0x9C, 0xA3, 0xAF)
GREY_SOFT   = RGBColor(0xD1, 0xD5, 0xDB)
SURFACE     = RGBColor(0xF9, 0xFA, 0xFB)
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = BLUE

HERE = Path(__file__).parent
ASSETS = HERE / "assets"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL = 9

# =========================================================================
# Drawing helpers
# =========================================================================
def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb, tf


def add_rect(slide, left, top, width, height, fill=BLUE, line=None,
             rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = 0.10
        except Exception:
            pass
    return shp


def add_circle(slide, cx, cy, radius, fill=BLUE, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - radius, cy - radius, radius * 2, radius * 2)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_top_bar(slide, color=BLUE, height=Inches(0.16)):
    add_rect(slide, Inches(0), Inches(0), SW, height, fill=color)


def add_footer(slide, slide_num, total=TOTAL):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             "Choolwe  •  The Case For Fossil Fuels",
             size=9, color=GREY_CTX)
    add_text(slide, Inches(11.8), Inches(7.15), Inches(1.2), Inches(0.3),
             f"{slide_num} / {total}",
             size=9, color=GREY_CTX, align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def section_tag(slide, label, color=BLUE, x=Inches(0.6), y=Inches(0.5)):
    add_circle(slide, x + Inches(0.10), y + Inches(0.16),
               Inches(0.07), fill=color)
    add_text(slide, x + Inches(0.30), y, Inches(8), Inches(0.4),
             label, size=11, bold=True, color=color)


# ---- Tiny icon compositions (native shapes, no clipart) ----
def icon_flame(slide, cx, cy, r=Inches(0.5), color=AMBER):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r, cy - r * 1.1, r * 2, r * 2.2)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r * 0.45, cy - r * 0.4, r * 0.9, r * 1.2)
    inner.fill.solid(); inner.fill.fore_color.rgb = AMBER_LIGHT
    inner.line.fill.background(); inner.shadow.inherit = False


def icon_bolt(slide, cx, cy, r=Inches(0.5), color=BLUE):
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    bolt = slide.shapes.add_shape(
        MSO_SHAPE.LIGHTNING_BOLT, cx - r * 0.5, cy - r * 0.75, r * 1.0, r * 1.5)
    bolt.fill.solid(); bolt.fill.fore_color.rgb = BG_WHITE
    bolt.line.fill.background(); bolt.shadow.inherit = False


def icon_globe(slide, cx, cy, r=Inches(0.5), color=BLUE):
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - r * 0.85, cy - r * 0.04, r * 1.7, r * 0.08)
    line.fill.solid(); line.fill.fore_color.rgb = BG_WHITE
    line.line.fill.background(); line.shadow.inherit = False
    mer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - r * 0.04, cy - r * 0.85, r * 0.08, r * 1.7)
    mer.fill.solid(); mer.fill.fore_color.rgb = BG_WHITE
    mer.line.fill.background(); mer.shadow.inherit = False


def icon_balance(slide, cx, cy, r=Inches(0.5), color=BLUE):
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - r * 0.7, cy - r * 0.05, r * 1.4, r * 0.10)
    bar.fill.solid(); bar.fill.fore_color.rgb = BG_WHITE
    bar.line.fill.background(); bar.shadow.inherit = False
    post = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - r * 0.05, cy - r * 0.05, r * 0.10, r * 0.65)
    post.fill.solid(); post.fill.fore_color.rgb = BG_WHITE
    post.line.fill.background(); post.shadow.inherit = False


# =========================================================================
# Slide 1 — Title (today's topic dominant, recap demoted)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), Inches(4.8), SH, fill=BLUE)
add_rect(s, Inches(4.8), Inches(0), Inches(0.16), SH, fill=AMBER)

add_text(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(0.5),
         "FOLLOW-UP", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.6), Inches(2.7), Inches(3.8), Inches(1.0),
         "The Case", size=46, bold=True, color=BG_WHITE)
add_text(s, Inches(0.6), Inches(3.45), Inches(3.8), Inches(1.0),
         "For Fossil", size=46, bold=True, color=BG_WHITE)
add_text(s, Inches(0.6), Inches(4.20), Inches(3.8), Inches(1.0),
         "Fuels", size=46, bold=True, color=BG_WHITE)
add_text(s, Inches(0.6), Inches(5.5), Inches(3.8), Inches(0.4),
         "A steelman counter-argument", size=14, color=BLUE_LIGHT)

# Right side — today's thesis leads; recap reduced to a kicker
add_text(s, Inches(5.4), Inches(2.0), Inches(7.5), Inches(0.4),
         "TODAY", size=11, bold=True, color=GREY_CTX)
add_text(s, Inches(5.4), Inches(2.5), Inches(7.5), Inches(1.6),
         "Keeping the lights —\nand the heat — on.",
         size=30, color=TEXT, bold=True)
add_text(s, Inches(5.4), Inches(4.2), Inches(7.5), Inches(0.9),
         "Why cutting fossil fuels too fast can cost more lives than it "
         "saves — in Lusaka and in Alberta alike.",
         size=14, color=TEXT_MUTED)
add_text(s, Inches(5.4), Inches(5.35), Inches(7.5), Inches(0.5),
         "Choolwe", size=22, bold=True, color=BLUE)
add_text(s, Inches(5.4), Inches(5.9), Inches(7.5), Inches(0.4),
         "Analytics in Training  •  May 2026", size=13, color=TEXT_MUTED)

icon_flame(s, Inches(11.7), Inches(6.7), r=Inches(0.32))
icon_bolt(s, Inches(12.5), Inches(6.7), r=Inches(0.32))

set_notes(s,
    "(~15 sec) Yesterday I argued we should cut fossil fuels. You all saw "
    "that talk, so I won't recap it. Today I want to make the most honest "
    "counter-argument I can — not because I've changed my mind on climate "
    "change, but because the picture I painted was incomplete. The story "
    "today is simple: reliable energy is survival, and clean energy alone "
    "can't guarantee it yet — not in Lusaka, and not in Alberta.")

# =========================================================================
# Slide 2 — Primer: what are fossil fuels
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s)
section_tag(s, "PRIMER  •  WHAT WE'RE TALKING ABOUT",
            color=BLUE, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.2),
         "Fossil fuels are stored ancient sunlight —\n"
         "and still about 80% of the energy the world runs on.",
         size=24, bold=True, color=TEXT)

primer_cards = [
    (icon_globe, "Coal",
     "Ancient land plants and swamp forests,\nburied and compressed. Mined as a solid.",
     BLUE),
    (icon_flame, "Oil (petroleum)",
     "Ancient ocean plankton and algae,\ncooked under pressure. Drilled as a liquid.",
     AMBER),
    (icon_bolt, "Natural gas",
     "The same organic matter, cooked hotter.\nDrilled as a gas — mostly methane.",
     BLUE),
]
card_w = Inches(3.85); card_h = Inches(3.0); gap = Inches(0.3)
start_x = Inches(0.6); y = Inches(2.9)
for i, (icon_fn, head, body, c) in enumerate(primer_cards):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, card_h, fill=SURFACE, rounded=True)
    add_rect(s, x, y, Inches(0.18), card_h, fill=c)
    icon_fn(s, x + Inches(0.7), y + Inches(0.65), r=Inches(0.32))
    add_text(s, x + Inches(1.25), y + Inches(0.45),
             card_w - Inches(1.4), Inches(0.8),
             head, size=16, bold=True, color=TEXT)
    add_text(s, x + Inches(0.4), y + Inches(1.55),
             card_w - Inches(0.6), Inches(1.3),
             body, size=12.5, color=TEXT_MUTED)

add_text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6),
         "What we use them for today:  electricity · transport · heat · "
         "industry (steel, cement, fertilizer) · cooking.",
         size=13, bold=True, color=TEXT_MUTED)

add_footer(s, 2)
set_notes(s,
    "(~25 sec) Quick grounding before the argument. Fossil fuels are just "
    "stored ancient sunlight — plants and ocean plankton, buried and cooked "
    "over millions of years into coal, oil and gas. Burning them releases "
    "that stored energy, and the CO₂ that comes with it. They're still "
    "roughly 80% of all the energy humanity uses — not just electricity, but "
    "transport, heat, and the steel, cement and fertilizer modern life is "
    "built from. Hold that scale in mind for everything that follows.")

# =========================================================================
# Slide 3 — Big Idea (reliability spine)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s)
add_rect(s, Inches(0), Inches(0.16), Inches(0.45), SH - Inches(0.16), fill=AMBER)
section_tag(s, "THE BIG IDEA", color=AMBER, x=Inches(0.9), y=Inches(0.7))

tb = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = Emu(0)
p = tf.paragraphs[0]
def _run(p, text, color, bold=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Calibri"
_run(p, "Reliable energy is survival — and ", TEXT_MUTED)
_run(p, "clean energy alone can't deliver it", AMBER, True)
_run(p, " when we need it most. From ", TEXT_MUTED)
_run(p, "Zambia's drought", AMBER, True)
_run(p, " to ", TEXT_MUTED)
_run(p, "a Canadian winter", AMBER, True)
_run(p, ", cutting fossil fuels now costs lives.", TEXT_MUTED)

add_text(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.0),
         "The honest path is a managed transition,\nnot abandonment.",
         size=40, bold=True, color=BLUE)

icon_balance(s, Inches(12.4), Inches(1.05), r=Inches(0.4))

add_footer(s, 3)
set_notes(s,
    "(~20 sec) This is what I want you to walk out of the room remembering. "
    "Reliable energy is survival — and the clean-energy-only plan can't "
    "guarantee it yet. We have lived proof on two continents: Zambia's "
    "drought and a Canadian winter. So the honest answer isn't abandonment — "
    "it's a managed transition. Let me show you how I got there.")

# =========================================================================
# Slide 4 — Complication: reliability test in the South (Zambia)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s, color=AMBER)
section_tag(s, "COMPLICATION  •  THE RELIABILITY TEST, IN THE SOUTH",
            color=AMBER, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4),
         "Zambia bet 85% of its grid on \"clean\" hydropower —\n"
         "and a climate-driven drought left us in the dark for 21 hours a day.",
         size=22, bold=True, color=TEXT)

s.shapes.add_picture(str(ASSETS / "zambia_mix.png"),
                     Inches(0.4), Inches(2.6), width=Inches(7.5))

card_x = Inches(8.4); card_w = Inches(4.5)
add_rect(s, card_x, Inches(2.7), card_w, Inches(1.7), fill=AMBER_LIGHT, rounded=True)
add_text(s, card_x + Inches(0.3), Inches(2.85), card_w - Inches(0.6), Inches(0.4),
         "THE BET", size=11, bold=True, color=AMBER)
add_text(s, card_x + Inches(0.3), Inches(3.2), card_w - Inches(0.6), Inches(1.2),
         "85% of installed capacity\nis hydropower.", size=18, bold=True, color=TEXT)

add_rect(s, card_x, Inches(4.55), card_w, Inches(1.7), fill=BLUE_LIGHT, rounded=True)
add_text(s, card_x + Inches(0.3), Inches(4.70), card_w - Inches(0.6), Inches(0.4),
         "THE LESSON", size=11, bold=True, color=BLUE)
add_text(s, card_x + Inches(0.3), Inches(5.05), card_w - Inches(0.6), Inches(1.2),
         "Firm, dispatchable power matters.\n"
         "Pretending otherwise is paid for in dark hospitals.",
         size=13, color=TEXT)

add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4),
         "The answer isn't \"burn more coal.\" It's \"don't bet the grid on one weather pattern.\"",
         size=13, color=TEXT_MUTED)

add_footer(s, 4)
set_notes(s,
    "(~40 sec) We lived this. We did the 'clean' thing — 85% of our installed "
    "capacity is hydropower. Then the climate broke our climate solution. The "
    "2024 drought collapsed reservoir levels and ZESCO ran load shedding of up "
    "to 21 hours a day. The lesson is NOT 'burn more coal.' The lesson is: firm, "
    "dispatchable power matters, and pretending otherwise is paid for in dark "
    "hospitals and empty cold-chains.")

# =========================================================================
# Slide 5 — Complication: reliability test in the North (Alberta cold snap)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s, color=AMBER)
section_tag(s, "COMPLICATION  •  THE RELIABILITY TEST, IN THE NORTH",
            color=AMBER, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4),
         "When it's coldest, clean energy is weakest —\n"
         "and fossil fuels are what keep people alive.",
         size=22, bold=True, color=TEXT)

s.shapes.add_picture(str(ASSETS / "cold_climate.png"),
                     Inches(0.4), Inches(2.6), width=Inches(7.6))

card_x = Inches(8.5); card_w = Inches(4.4)
add_rect(s, card_x, Inches(2.7), card_w, Inches(1.85), fill=AMBER_LIGHT, rounded=True)
add_text(s, card_x + Inches(0.3), Inches(2.85), card_w - Inches(0.6), Inches(0.4),
         "WHY WARMTH IS NON-NEGOTIABLE", size=11, bold=True, color=AMBER)
add_text(s, card_x + Inches(0.3), Inches(3.2), card_w - Inches(0.6), Inches(0.7),
         "~9× more", size=30, bold=True, color=AMBER)
add_text(s, card_x + Inches(0.3), Inches(3.95), card_w - Inches(0.6), Inches(0.6),
         "deaths from cold than from heat, worldwide.", size=13, color=TEXT)

add_rect(s, card_x, Inches(4.7), card_w, Inches(1.55), fill=BLUE_LIGHT, rounded=True)
add_text(s, card_x + Inches(0.3), Inches(4.85), card_w - Inches(0.6), Inches(0.4),
         "GAS HEATS HALF OF CANADIAN HOMES", size=11, bold=True, color=BLUE)
add_text(s, card_x + Inches(0.3), Inches(5.2), card_w - Inches(0.6), Inches(1.0),
         "72% in Alberta, 62% in Ontario. In a -30°C winter, "
         "\"no fossil fuels\" means no heat.",
         size=13, color=TEXT)

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "It isn't only the poor South. A rich, modern province nearly went "
         "dark — and only firm power kept it alive.",
         size=13, color=TEXT_MUTED)

add_footer(s, 5)
set_notes(s,
    "(~45 sec) And it isn't only the poor South. In January 2024, Alberta — a "
    "rich, modern province — hit a minus-forty cold snap. Demand smashed a "
    "record at 12,384 megawatts. The province has over six thousand megawatts "
    "of wind and solar, but at the evening peak the sun had set and the wind "
    "was calm, so it delivered almost nothing. The grid fell to TEN megawatts "
    "of reserve and the government sent an emergency text telling everyone to "
    "cut power to essentials. Gas-fired plants carried the province. And "
    "remember — worldwide, cold kills about nine times as many people as heat. "
    "In most of the cold world, staying warm runs on fossil fuels: gas heats "
    "roughly half of Canadian homes. 'No fossil fuels' in a minus-thirty "
    "winter means no heat.")

# =========================================================================
# Slide 6 — Complication: justice asymmetry (isotype)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s, color=AMBER)
section_tag(s, "COMPLICATION  •  THE JUSTICE ASYMMETRY",
            color=AMBER, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4),
         "Africa hosts 17% of the world's people, caused under 3% of cumulative\n"
         "emissions — and is being asked to skip the ladder every wealthy country climbed.",
         size=21, bold=True, color=TEXT)

s.shapes.add_picture(str(ASSETS / "africa_asymmetry.png"),
                     Inches(0.4), Inches(2.9), width=Inches(12.5))

add_rect(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.95),
         fill=BLUE_LIGHT, rounded=True)
add_text(s, Inches(0.95), Inches(6.10), Inches(11.5), Inches(0.4),
         "Every wealthy country got rich by burning coal, then oil, then gas.",
         size=14, bold=True, color=BLUE)
add_text(s, Inches(0.95), Inches(6.45), Inches(11.5), Inches(0.4),
         "Even the clean-energy supply chain itself — steel, cement, shipping — still runs on fossil fuels today.",
         size=12, color=TEXT_MUTED)

add_footer(s, 6)
set_notes(s,
    "(~40 sec) Now the justice layer. Every wealthy country on Earth got rich "
    "by burning coal, then oil, then gas. They are now asking the continent that "
    "contributed least to the problem to skip that ladder — using clean-energy "
    "technology whose OWN supply chain still runs on fossil fuels. That's not "
    "a climate plan. That's an injustice with a green sticker on it.")

# =========================================================================
# Slide 7 — Bridge: energy & life expectancy
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s, color=AMBER)
section_tag(s, "BRIDGE  •  WHAT ENERGY ACCESS BUYS",
            color=AMBER, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4),
         "Below ~3,000 kWh per person per year, every extra kWh\n"
         "buys life-years. Zambia sits right on the steep part of the curve.",
         size=22, bold=True, color=TEXT)

s.shapes.add_picture(str(ASSETS / "energy_life.png"),
                     Inches(0.6), Inches(2.8), width=Inches(12.1))

add_footer(s, 7)
set_notes(s,
    "(~30 sec) This is the Hans Rosling curve. Below about three thousand "
    "kilowatt-hours per person per year, every additional kilowatt-hour "
    "translates directly into life-years — through clinics, refrigerated "
    "vaccines, lit classrooms, water pumps. Zambia is right on the steep "
    "part of that curve. Cutting energy access there costs life-years "
    "immediately. That's why the moral arithmetic doesn't end with climate.")

# =========================================================================
# Slide 8 — Resolution
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s, color=BLUE)
section_tag(s, "RESOLUTION  •  THE MANAGED TRANSITION",
            color=BLUE, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4),
         "Gas as a bridge, renewables stacked on top, firm power kept\n"
         "until storage can carry the load — in that order.",
         size=22, bold=True, color=TEXT)

block_w = Inches(3.95); block_h = Inches(3.6)
ys = Inches(3.0)
xs = [Inches(0.6), Inches(4.7), Inches(8.8)]
blocks = [
    ("NOW — BRIDGE", "Natural gas",
     "Emits 50–60% less CO₂ than coal.\n\n"
     "Firm and dispatchable — the power that\n"
     "carried Alberta and would have spared\n"
     "Zambia.\n\n"
     "Africa holds ~13% of proven gas reserves.",
     BLUE, BG_WHITE, True, icon_flame),
    ("BUILD-OUT", "Renewables, hard",
     "Solar in Zambia is already the\n"
     "cheapest new capacity.\n\n"
     "Wind, hydro diversification,\n"
     "storage — scaled in parallel,\n"
     "not waited for.",
     SURFACE, TEXT, False, icon_bolt),
    ("THEN — RETIRE", "Phase fossil fuels down",
     "When grids are firm and storage can\n"
     "carry the load, retire gas plants.\n\n"
     "Not before. Not as a slogan.\n"
     "As an engineering schedule.",
     SURFACE, TEXT, False, icon_globe),
]
for x, (tag, head, body, fill, fg, is_lead, icon_fn) in zip(xs, blocks):
    add_rect(s, x, ys, block_w, block_h, fill=fill, rounded=True)
    if is_lead:
        icon_fn(s, x + block_w - Inches(0.55), ys + Inches(0.55),
                r=Inches(0.35), color=BG_WHITE)
    else:
        icon_fn(s, x + block_w - Inches(0.55), ys + Inches(0.55), r=Inches(0.30))
    tag_color = BG_WHITE if is_lead else BLUE
    add_text(s, x + Inches(0.35), ys + Inches(0.35),
             block_w - Inches(1.4), Inches(0.4),
             tag, size=11, bold=True, color=tag_color)
    add_text(s, x + Inches(0.35), ys + Inches(0.85),
             block_w - Inches(0.6), Inches(0.7),
             head, size=20, bold=True, color=fg)
    add_text(s, x + Inches(0.35), ys + Inches(1.8),
             block_w - Inches(0.6), block_h - Inches(2.0),
             body, size=12, color=(BG_WHITE if is_lead else TEXT_MUTED))

for i in range(2):
    ax_ = xs[i] + block_w + Inches(0.0)
    arrow = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, ax_ - Inches(0.05),
        ys + Inches(1.55), Inches(0.30), Inches(0.45))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = GREY_CTX
    arrow.line.fill.background(); arrow.shadow.inherit = False

add_footer(s, 8)
set_notes(s,
    "(~50 sec) I'm not arguing fossil fuels forever. I'm arguing for a MANAGED "
    "transition. Natural gas as a bridge — emits half what coal does, and it's "
    "the firm power that carried Alberta and would have spared Zambia; Africa "
    "holds 13% of global reserves. Renewables scaled aggressively on top — solar "
    "in Zambia is already the cheapest new capacity. And firm power kept on until "
    "storage and grid infrastructure can stand alone. That's how we keep the "
    "lights and the heat on while decarbonising.")

# =========================================================================
# Slide 9 — Sources
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_top_bar(s, color=BLUE)
section_tag(s, "APPENDIX  •  SOURCES", color=BLUE, x=Inches(0.6), y=Inches(0.5))

add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
         "Every statistic on every slide, with a citation.",
         size=22, bold=True, color=TEXT)
add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
         "Full notes preserved in the repo at  research/notes.md.",
         size=12, color=TEXT_MUTED)

sources = [
    ("Africa electricity access (~600M without)",
     "IEA, Financing Electricity Access in Africa (2025)"),
    ("Africa cumulative CO₂ emissions (<3%; sub-Saharan ~0.55%)",
     "Energy for Growth Hub (citing Our World in Data)"),
    ("Africa population share (~17% global)",
     "UN World Population Prospects"),
    ("Zambia capacity mix (85% hydropower; ~3,986 MW)",
     "Global Legal Insights, Energy Laws and Regulations 2025: Zambia"),
    ("Zambia 2024 load shedding (up to 21 hrs/day)",
     "Pulitzer Center; Afrobarometer AD1178 (2025)"),
    ("Cold kills ~9× more than heat (~4.6M vs ~489k deaths/yr)",
     "The Lancet Planetary Health (Gasparrini et al.)"),
    ("Canadian homes heated by gas (~half; 72% AB, 62% ON)",
     "Statistics Canada — How Canadians heat their home"),
    ("Alberta cold-snap grid alert (record 12,384 MW; reserves to 10 MW)",
     "AESO / Alberta MSA, Jan & Apr 2024 Event Report"),
    ("Natural gas vs coal CO₂ (50–60% less); Africa gas reserves (~13%)",
     "Atlantic Council, Natural gas in Africa's energy transition"),
    ("Steel & cement embodied emissions in renewables",
     "MIT Climate Portal"),
    ("Life expectancy vs electricity per capita (illustrative)",
     "World Bank Indicators (kWh/capita; life expectancy at birth)"),
]
col_x = [Inches(0.6), Inches(6.95)]
col_w = Inches(6.0)
row_y = Inches(2.4)
row_h = Inches(0.40)
for i, (claim, src) in enumerate(sources):
    col = i % 2
    row = i // 2
    x = col_x[col]; y = row_y + row * (row_h + Inches(0.20))
    add_circle(s, x + Inches(0.1), y + Inches(0.15), Inches(0.05), fill=AMBER)
    add_text(s, x + Inches(0.25), y, col_w - Inches(0.25), Inches(0.25),
             claim, size=11, bold=True, color=TEXT)
    add_text(s, x + Inches(0.25), y + Inches(0.22),
             col_w - Inches(0.25), Inches(0.25),
             src, size=10, color=TEXT_MUTED)

add_footer(s, 9)
set_notes(s,
    "Reference slide for Q&A. If anyone questions a stat, the source is here. "
    "Full URLs are in the repo's research/notes.md.")

# =========================================================================
# Save (v2 names only — originals untouched)
# =========================================================================
out = HERE / "deliverable" / "case-for-fossil-fuels-v2.pptx"
out.parent.mkdir(exist_ok=True)
prs.save(out)
print(f"wrote {out}")
